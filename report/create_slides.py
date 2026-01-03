"""
Generate PowerPoint slides for:
"Revisiting DeepLabV3: Efficient Context and Attention Enhancements for Semantic Segmentation"
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0, 51, 102)
ACCENT_BLUE = RGBColor(0, 112, 192)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GREEN = RGBColor(0, 128, 0)
RED = RGBColor(192, 0, 0)


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, bullet_points, sub_bullets=None):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_after = Pt(12)
        p.level = 0
        
        # Add sub-bullets if provided
        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                p = tf.add_paragraph()
                p.text = "    ‣ " + sub
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(80, 80, 80)
                p.space_after = Pt(6)
    
    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table_width = Inches(12) if col_widths is None else sum(col_widths)
    left = (prs.slide_width - table_width) / 2
    
    table = slide.shapes.add_table(num_rows, num_cols, left, Inches(1.6), table_width, Inches(0.5 * num_rows)).table
    
    # Set column widths
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = width
    
    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    return slide


def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    return slide


def add_section_slide(prs, section_title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide


# ============== CREATE SLIDES ==============

# Slide 1: Title
add_title_slide(
    prs,
    "Revisiting DeepLabV3: Efficient Context and\nAttention Enhancements for Semantic Segmentation",
    "Lam T. Le, Thong C. Nguyen, Thong H. Nguyen\nVNU-HCM University of Science"
)

# Slide 2: Motivation
add_content_slide(
    prs,
    "Motivation",
    [
        "Semantic segmentation: Assign class label to every pixel",
        "Key challenge: Balance multi-scale context with spatial precision",
        "CNN-based methods remain practical for deployment:",
        "Research question: Can lightweight attention improve DeepLabV3+ without sacrificing efficiency?"
    ],
    sub_bullets={
        2: ["Mobile and edge devices", "Resource-constrained scenarios", "Real-time applications"]
    }
)

# Slide 3: DeepLabV3+ Background
add_content_slide(
    prs,
    "Background: DeepLabV3+ Architecture",
    [
        "Encoder-Decoder framework for semantic segmentation",
        "ASPP (Atrous Spatial Pyramid Pooling) module:",
        "Decoder fuses low-level (48ch) + high-level (256ch) features",
        "Backbone options: ResNet, MobileNetV2, MobileNetV3"
    ],
    sub_bullets={
        1: ["1×1 convolution", "3 atrous convolutions (rates: 6, 12, 18)", "Global Average Pooling", "Concatenate → Project to 256 channels"]
    }
)

# Slide 4: Attention Mechanisms
add_table_slide(
    prs,
    "Attention Mechanisms Overview",
    ["Module", "Key Idea", "Overhead"],
    [
        ["Shuffle Attention (SA)", "Group-wise channel + spatial attention", "Negligible"],
        ["ECA", "1D conv on channel descriptors", "O(k) params"],
        ["EPSA", "Multi-scale splits + SE-style attention", "Moderate"],
        ["Strip Pooling", "Directional (H/V) long-range pooling", "Moderate"]
    ],
    col_widths=[Inches(3), Inches(5.5), Inches(2.5)]
)

# Slide 5: SA-ECA Approach
add_content_slide(
    prs,
    "Proposed Approach 1: SA-ECA",
    [
        "Shuffle Attention after ASPP concatenation (1280 channels)",
        "ECA at decoder fusion (304 channels)",
        "Forward pass:",
        "Design rationale: Lightweight recalibration at both encoder and decoder"
    ],
    sub_bullets={
        2: [
            "ASPP branches → Concat → ShuffleAttention → 1×1 Conv",
            "Upsample → Concat with low-level features → ECA → Classifier"
        ]
    }
)

# Slide 6: SP-EPSA Approach
add_content_slide(
    prs,
    "Proposed Approach 2: SP-EPSA",
    [
        "Replace Global Average Pooling with Strip Pooling in ASPP",
        "Add EPSA after ASPP projection (256 channels)",
        "Apply ECA on low-level features before fusion",
        "Design rationale: Capture directional context for structured scenes"
    ],
    sub_bullets={
        0: ["Horizontal pooling: captures road-like structures", "Vertical pooling: captures buildings, poles"]
    }
)

# Slide 7: Experimental Setup
add_table_slide(
    prs,
    "Experimental Setup",
    ["Setting", "PASCAL VOC 2012", "Cityscapes"],
    [
        ["Classes", "21", "19"],
        ["Training Images", "10,582 (augmented)", "2,975"],
        ["Validation Images", "1,449", "500"],
        ["Crop Size", "513 × 513", "768 × 768"],
        ["Total Iterations", "30,000", "30,000"],
        ["LR Policy", "Poly (p=0.9)", "Poly (p=0.9)"],
        ["Base LR", "0.01", "0.1"],
        ["Backbone", "MobileNetV2 / MNv3-L", "MobileNetV2 / MNv3-L"]
    ],
    col_widths=[Inches(3.5), Inches(3.5), Inches(3.5)]
)

# Slide 8: Model Complexity
add_table_slide(
    prs,
    "Model Complexity Comparison",
    ["Model", "Backbone", "Params (M)", "GFLOPs", "VOC mIoU", "City mIoU"],
    [
        ["Baseline", "MNv2", "5.23", "17.01", "66.65%", "72.07%"],
        ["Baseline", "MNv3-L", "11.14", "21.69", "54.23%", "63.74%"],
        ["SA-ECA", "MNv2", "5.23", "17.01", "67.96% ↑", "71.05%"],
        ["SP-EPSA", "MNv2", "5.93", "17.63", "65.45%", "74.62% ↑"]
    ],
    col_widths=[Inches(2), Inches(1.8), Inches(2), Inches(2), Inches(2), Inches(2)]
)

# Slide 9: VOC Results
add_content_slide(
    prs,
    "Results: PASCAL VOC 2012",
    [
        "Best: SA-ECA with MobileNetV2 → 67.96% mIoU (+1.31%)",
        "Improvements across all metrics (OA, MA, FWAcc)",
        "SP-EPSA shows marginal decrease (-1.2% mIoU)",
        "MobileNetV3-Large consistently underperforms MobileNetV2"
    ],
    sub_bullets={
        0: ["Effective for diverse object categories", "Negligible computational overhead"],
        2: ["Strip pooling less suited for irregular object layouts"]
    }
)

# Slide 10: Cityscapes Results
add_content_slide(
    prs,
    "Results: Cityscapes",
    [
        "Best: SP-EPSA with MobileNetV2 → 74.62% mIoU (+2.55%)",
        "Strip Pooling excels at capturing urban structures:",
        "SA-ECA slightly underperforms baseline (-1.02% mIoU)",
        "Demonstrates dataset-dependent attention effectiveness"
    ],
    sub_bullets={
        1: ["Horizontal structures: roads, sidewalks", "Vertical structures: buildings, poles, traffic signs"]
    }
)

# Slide 11: Key Findings - Dataset Dependence
add_table_slide(
    prs,
    "Key Finding: Dataset-Dependent Behavior",
    ["Dataset", "Best Variant", "Improvement", "Reason"],
    [
        ["PASCAL VOC", "SA-ECA", "+1.31%", "Diverse objects, irregular layouts"],
        ["Cityscapes", "SP-EPSA", "+2.55%", "Structured scenes, elongated elements"]
    ],
    col_widths=[Inches(2.5), Inches(2.5), Inches(2), Inches(4)]
)

# Slide 12: Key Findings - Backbone
add_content_slide(
    prs,
    "Key Finding: Backbone Interactions",
    [
        "MobileNetV3-Large consistently underperforms MobileNetV2",
        "Hypothesis: Built-in SE blocks conflict with external attention",
        "MobileNetV2's simpler architecture → better substrate for augmentation",
        "Training instability observed with MobileNetV3-Large + attention"
    ],
    sub_bullets={
        0: ["Despite having 2× more parameters (11.14M vs 5.23M)"],
        3: ["Oscillating validation metrics", "Suboptimal convergence even with warmup"]
    }
)

# Slide 13: Computational Efficiency
add_two_column_slide(
    prs,
    "Computational Efficiency Analysis",
    "SA-ECA Variant",
    [
        "Negligible parameter increase",
        "Same GFLOPs as baseline",
        "Shuffle Attention: group-wise ops",
        "ECA: single 1D conv (k ≤ 7)",
        "Ideal for strict resource limits"
    ],
    "SP-EPSA Variant",
    [
        "+0.7M parameters (+13%)",
        "+0.62 GFLOPs (+3.6%)",
        "Strip Pooling: 2 pooling ops",
        "EPSA: depthwise separable convs",
        "Better accuracy-efficiency trade-off"
    ]
)

# Slide 14: Limitations
add_two_column_slide(
    prs,
    "Limitations & Future Work",
    "Limitations",
    [
        "Validation-only evaluation",
        "No component-wise ablations",
        "Hyperparameters not extensively tuned",
        "MobileNetV3-L integration challenges",
        "Bundled changes in SP-EPSA"
    ],
    "Future Work",
    [
        "Ablation study: isolate each module's contribution",
        "Test server evaluation for Cityscapes",
        "Investigate SE-block interference",
        "Extend to other datasets",
        "Explore attention placement strategies"
    ]
)

# Slide 15: Conclusion
add_content_slide(
    prs,
    "Conclusion",
    [
        "Proposed two lightweight attention variants for DeepLabV3+:",
        "Minimal overhead: +13% params, +3.6% GFLOPs for SP-EPSA",
        "Key insight: Attention effectiveness is dataset-dependent",
        "Simpler backbones (MobileNetV2) better suited for attention augmentation",
        "Practical contribution: Efficient enhancements for mobile deployment"
    ],
    sub_bullets={
        0: ["SA-ECA: Best for diverse datasets (VOC: +1.31%)", "SP-EPSA: Best for structured scenes (Cityscapes: +2.55%)"]
    }
)

# Slide 16: Thank You
add_title_slide(
    prs,
    "Thank You",
    "Questions?\n\nContact: ltlam22@apcs.fitus.edu.vn"
)

# Save the presentation
output_path = "DeepLabV3_Attention_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
